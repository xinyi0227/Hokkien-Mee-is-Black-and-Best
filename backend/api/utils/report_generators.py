from reportlab.lib.pagesizes import A4, letter
from reportlab.platypus import (
    SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image as RLImage, PageBreak
)
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib import colors
from reportlab.pdfgen import canvas

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

from io import BytesIO
import datetime, requests
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import seaborn as sns
import pandas as pd
import numpy as np


# ---------------------------
# Helpers: styles and branding
# ---------------------------

PRIMARY_BLUE = colors.HexColor("#143D6B")
ACCENT_GREEN = colors.HexColor("#1E9E6A")
MUTED_TEXT = colors.HexColor("#4A5568")
LIGHT_BG = colors.HexColor("#F7FAFC")
TABLE_HEADER_BG = colors.HexColor("#2D3748")
TABLE_ROW_ALT = colors.HexColor("#EDF2F7")
RED = colors.HexColor("#C53030")
ORANGE = colors.HexColor("#DD6B20")

def _pdf_header_footer(canvas_obj, doc):
    canvas_obj.setFillColor(MUTED_TEXT)
    canvas_obj.setFont("Helvetica", 8)
    text = f"AI Business Analytics • Generated {datetime.datetime.now().strftime('%b %d, %Y %I:%M %p')}"
    canvas_obj.drawRightString(doc.pagesize[0] - 40, 20, text)
    canvas_obj.setFillColor(PRIMARY_BLUE)
    canvas_obj.setLineWidth(0.3)
    canvas_obj.line(40, doc.pagesize[1]-40, doc.pagesize[0]-40, doc.pagesize[1]-40)
    canvas_obj.setFillColor(MUTED_TEXT)
    canvas_obj.setFont("Helvetica-Bold", 9)
    canvas_obj.drawString(40, doc.pagesize[1]-30, "Executive Business Intelligence Report")

def _styles():
    styles = getSampleStyleSheet()
    title = ParagraphStyle(
        'TitleH1',
        parent=styles['Heading1'],
        fontSize=22,
        leading=26,
        spaceAfter=18,
        textColor=PRIMARY_BLUE,
        alignment=1,
        fontName='Helvetica-Bold'
    )
    h2 = ParagraphStyle(
        'H2',
        parent=styles['Heading2'],
        fontSize=16,
        leading=20,
        spaceBefore=12,
        spaceAfter=10,
        textColor=PRIMARY_BLUE,
        fontName='Helvetica-Bold'
    )
    h3 = ParagraphStyle(
        'H3',
        parent=styles['Heading3'],
        fontSize=12.5,
        leading=16,
        spaceBefore=8,
        spaceAfter=6,
        textColor=colors.HexColor("#9B2C2C"),
        fontName='Helvetica-Bold'
    )
    normal = ParagraphStyle(
        'Body',
        parent=styles['Normal'],
        fontSize=10.5,
        leading=14,
        textColor=MUTED_TEXT
    )
    bullet = ParagraphStyle(
        'Bullet',
        parent=styles['Normal'],
        fontSize=10.5,
        leading=14,
        leftIndent=16,
        bulletIndent=8,
        textColor=MUTED_TEXT
    )
    note = ParagraphStyle(
        'Note',
        parent=styles['Italic'],
        fontSize=9.5,
        textColor=colors.HexColor("#2F855A")
    )
    return title, h2, h3, normal, bullet, note


# ---------------------------
# Cleaning Report
# ---------------------------
from reportlab.platypus import Image as RLImage
import matplotlib
from reportlab.platypus import Image
matplotlib.use('Agg')  # Use non-interactive backend

class CleaningReportGenerator:
    def create_cleaning_report(self, cleaning_log, filename):
        buffer = BytesIO()
        doc = SimpleDocTemplate(
            buffer, pagesize=A4, topMargin=0.6*inch, bottomMargin=0.6*inch, leftMargin=0.7*inch, rightMargin=0.7*inch
        )
        title, h2, h3, normal, bullet, note = _styles()
        story = []

        # Cover/title block
        story.append(Paragraph(f"Data Cleaning Report", title))
        story.append(Paragraph(f"Dataset: {filename}", normal))
        story.append(Spacer(1, 10))
        story.append(Paragraph("This report summarizes detected issues, corrective actions, and key quality metrics from automated data cleaning.", note))
        story.append(Spacer(1, 16))

        # Summary table
        story.append(Paragraph("Cleaning Summary", h2))
        summary_data = [
            ["Metric", "Value"],
            ["Original Rows", str(cleaning_log['original_shape'][0])],  # Access first element
            ["Original Columns", str(cleaning_log['original_shape'][1])], # Access second element
            ["Final Rows", str(cleaning_log['final_shape'][0])],         # Access first element
            ["Final Columns", str(cleaning_log['final_shape'][1])],      # Access second element
            ["Total Issues Found", str(cleaning_log['summary']['total_issues_found'])],
            ["Total Actions Taken", str(cleaning_log['summary']['total_actions_taken'])],
        ]
        tbl = Table(summary_data, hAlign='LEFT', colWidths=[2.2*inch, 3.0*inch])
        tbl.setStyle(TableStyle([
            ('BACKGROUND', (0,0), (-1,0), TABLE_HEADER_BG),
            ('TEXTCOLOR', (0,0), (-1,0), colors.whitesmoke),
            ('FONTNAME', (0,0), (-1,0), 'Helvetica-Bold'),
            ('FONTSIZE', (0,0), (-1,0), 11),
            ('ALIGN', (0,0), (-1,0), 'LEFT'),
            ('ROWBACKGROUNDS', (0,1), (-1,-1), [LIGHT_BG, TABLE_ROW_ALT]),
            ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor("#CBD5E0")),
            ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
        ]))
        story.append(tbl)
        story.append(Spacer(1, 16))

        # Issues found
        if cleaning_log['issues_found']:
            story.append(Paragraph("Issues Identified", h2))
            for i, issue in enumerate(cleaning_log['issues_found'], 1):
                story.append(Paragraph(f"• {issue}", bullet))
            story.append(Spacer(1, 12))

        # Actions taken
        if cleaning_log['actions_taken']:
            story.append(Paragraph("Corrective Actions Applied", h2))
            for i, action in enumerate(cleaning_log['actions_taken'], 1):
                story.append(Paragraph(f"• {action}", bullet))
            story.append(Spacer(1, 12))

        # Recommendations
        story.append(Paragraph("Recommendations", h2))
        recs = [
            "Implement input validation at data entry points to reduce missing or invalid values.",
            "Schedule routine data quality checks and alerts for anomalies or schema drift.",
            "Adopt a standardized data dictionary and enforce consistent naming and typing.",
            "Automate recurring cleaning steps to ensure reproducibility and save analyst time."
        ]
        for r in recs:
            story.append(Paragraph(f"• {r}", bullet))

        doc.build(story, onFirstPage=_pdf_header_footer, onLaterPages=_pdf_header_footer)
        return buffer.getvalue()


# ---------------------------
# PDF Generator
# ---------------------------

class PDFGenerator:
    def create_analysis_report(self, gemini_output, filename, data_df=None, data_type=None):
        # If data provided, use the chart-embedded version
        if data_df is not None and isinstance(data_df, pd.DataFrame) and not data_df.empty:
            return self.create_analysis_report_with_charts(gemini_output, data_df, filename, data_type)
        # Otherwise fallback to text-only
        buffer = BytesIO()
        doc = SimpleDocTemplate(
            buffer, pagesize=A4, topMargin=0.6*inch, bottomMargin=0.6*inch, leftMargin=0.7*inch, rightMargin=0.7*inch
        )
        title, h2, h3, normal, bullet, note = _styles()
        story = []

        # Title block
        story.append(Paragraph(self._title_for_type(filename, data_type), title))
        story.append(Paragraph(self._subtitle_for_type(data_type), note))
        story.append(Spacer(1, 14))

        # Parse Gemini output
        if isinstance(gemini_output, dict) and 'full_report' in gemini_output:
            report_content = gemini_output['full_report']
            chart_recommendations = gemini_output.get('chart_recommendations', [])
        else:
            report_content = str(gemini_output)
            chart_recommendations = []

        # Convert report content
        self.process_gemini_report_content(story, report_content, title, h2, normal, bullet)

        # Chart recommendations (text)
        if chart_recommendations:
            story.append(Spacer(1, 10))
            story.append(Paragraph("Recommended Visualizations", h2))
            for chart_rec in chart_recommendations:
                clean_rec = chart_rec.replace('📊', '').replace('📈', '').replace('🎯', '').strip()
                story.append(Paragraph(f"• {clean_rec}", bullet))

        story.append(Spacer(1, 18))
        story.append(Paragraph("*Report generated by AI Business Analytics System*", note))
        story.append(Paragraph(f"*Generated: {datetime.datetime.now().strftime('%B %d, %Y %I:%M %p')}*", note))

        doc.build(story, onFirstPage=_pdf_header_footer, onLaterPages=_pdf_header_footer)
        return buffer.getvalue()

    def create_analysis_report_with_charts(self, gemini_output, data_df, filename, data_type=None):
        buffer = BytesIO()
        doc = SimpleDocTemplate(
            buffer, pagesize=A4, topMargin=0.6*inch, bottomMargin=0.6*inch, leftMargin=0.7*inch, rightMargin=0.7*inch
        )
        title, h2, h3, normal, bullet, note = _styles()
        story = []

        # Title block
        story.append(Paragraph(self._title_for_type(filename, data_type), title))
        story.append(Paragraph(self._subtitle_for_type(data_type), note))
        story.append(Spacer(1, 14))

        if isinstance(gemini_output, dict) and 'ai_insights' in gemini_output:
            report_content = gemini_output['ai_insights']
        elif isinstance(gemini_output, dict) and 'full_report' in gemini_output:
            report_content = gemini_output['full_report']
        else:
            report_content = str(gemini_output)

        # Content sections
        self.process_gemini_report_content(story, report_content, title, h2, normal, bullet)

        # Visualizations section
        story.append(PageBreak())
        story.append(Paragraph("Data Visualizations", h2))
        story.append(Spacer(1, 6))
        viz = self.generate_visuals_by_type(data_df, data_type)

        # Embed charts
        if viz:
            for caption, img_buf in viz:
                story.append(Paragraph(caption, h3))
                story.append(Spacer(1, 6))
                try:
                    story.append(RLImage(img_buf, width=6.5*inch, height=3.8*inch))
                except Exception:
                    # If size fails, try without constraining
                    story.append(RLImage(img_buf))
                story.append(Spacer(1, 16))
        else:
            story.append(Paragraph("No charts could be generated from the available data.", normal))

        story.append(Spacer(1, 10))
        story.append(Paragraph("*Report generated by AI Business Analytics System*", note))
        story.append(Paragraph(f"*Generated: {datetime.datetime.now().strftime('%B %d, %Y %I:%M %p')}*", note))

        doc.build(story, onFirstPage=_pdf_header_footer, onLaterPages=_pdf_header_footer)
        return buffer.getvalue()

    def create_specialized_analysis_report(self, analysis_results, filename, df_cleaned, data_type):
        """Create a specialized analysis report based on data type"""
        # This method can delegate to the existing chart-enabled method
        return self.create_analysis_report_with_charts(analysis_results, df_cleaned, filename, data_type)

    def _title_for_type(self, filename, data_type):
        base = "Business Analysis Report"
        if data_type == 'sales':
            base = "Sales Performance Report"
        elif data_type == 'financial':
            base = "Financial Performance Report"
        elif data_type == 'social_media':
            base = "Social Media Performance Report"
        return f"{base} – {filename}"

    def _subtitle_for_type(self, data_type):
        if data_type == 'sales':
            return "Focus: Revenue, product mix, volume trends, and customer behavior."
        if data_type == 'financial':
            return "Focus: Revenue, COGS, margins, OpEx, and profitability trends."
        if data_type == 'social_media':
            return "Focus: Engagement, platform effectiveness, content performance, and growth."
        return "Focus: Executive-ready insights, trends, and recommendations."

    def process_gemini_report_content(self, story, report_content, title_style, section_style, normal_style, bullet_style):
        # Robust Markdown-like parsing into paragraphs/bullets
        lines = report_content.split('\n')
        for raw in lines:
            line = raw.strip()
            if not line:
                continue
            if line.startswith('# '):
                story.append(Spacer(1, 6))
                story.append(Paragraph(line[2:].strip(), title_style))
                story.append(Spacer(1, 10))
            elif line.startswith('## '):
                story.append(Spacer(1, 6))
                story.append(Paragraph(line[3:].strip(), section_style))
            elif line.startswith('### '):
                sub_style = ParagraphStyle(
                    'SubH3', parent=normal_style, fontName='Helvetica-Bold',
                    fontSize=12, textColor=colors.HexColor("#9B2C2C"), spaceBefore=8, spaceAfter=6
                )
                story.append(Paragraph(line[4:].strip(), sub_style))
            elif any(emoji in line for emoji in ['📊', '📈', '🎯']):
                clean = line.replace('📊', '').replace('📈', '').replace('🎯', '').strip()
                callout = ParagraphStyle(
                    'Callout', parent=normal_style, textColor=ACCENT_GREEN
                )
                story.append(Paragraph(f"• {clean}", callout))
            elif line.startswith(('- ', '* ')):
                story.append(Paragraph(f"• {line[2:].strip()}", bullet_style))
            elif line[:3] in [f"{i}. " for i in range(1,10)]:
                story.append(Paragraph(line, bullet_style))
            else:
                story.append(Paragraph(line, normal_style))

    # ---------- Visualization packs ----------
    def generate_visuals_by_type(self, df, data_type):
        # Set style
        try:
            sns.set_style("whitegrid")
            plt.style.use('seaborn-v0_8-whitegrid')
        except Exception:
            pass

        if data_type == 'financial':
            return self._financial_charts(df)
        elif data_type == 'sales':
            return self._sales_charts(df)
        elif data_type == 'social_media':
            return self._social_media_charts(df)
        else:
            return self._general_charts(df)

    def _buf_from_fig(self, fig):
        buf = BytesIO()
        fig.savefig(buf, format='png', dpi=300, bbox_inches='tight', facecolor='white')
        buf.seek(0)
        plt.close(fig)
        return buf

    def _first_date(self, df):
        candidates = [c for c in df.columns if 'date' in c.lower() or 'time' in c.lower()]
        if not candidates:
            return None
        c = candidates[0]
        try:
            ser = pd.to_datetime(df[c], errors='coerce')
            return c, ser
        except Exception:
            return None

    # Financial charts: revenue/profit trend, margin trend, cost breakdown
    def _financial_charts(self, df):
        visuals = []
        cols = {c.lower(): c for c in df.columns}
        rev = next((cols[c] for c in cols if 'revenue' in c or ('total_revenue' in c)), None)
        profit = next((cols[c] for c in cols if 'net_profit' in c or c == 'profit' or 'gross_profit' in c), None)
        margin = next((cols[c] for c in cols if 'margin' in c), None)

        # 1) Revenue and Profit Trend
        try:
            date_info = self._first_date(df)
            if date_info and rev and profit:
                dc, ser = date_info
                tmp = df.copy()
                tmp[dc] = pd.to_datetime(tmp[dc], errors='coerce')
                tmp = tmp.dropna(subset=[dc])
                agg = tmp.groupby(tmp[dc].dt.to_period('D'))[[rev, profit]].sum()
                fig, ax1 = plt.subplots(figsize=(10,5.5))
                ax1.plot(agg.index.astype(str), agg[rev], color="#2E86AB", marker='o', label='Revenue')
                ax1.set_ylabel('Revenue', color="#2E86AB")
                ax2 = ax1.twinx()
                ax2.plot(agg.index.astype(str), agg[profit], color="#C0392B", marker='s', label='Profit')
                ax2.set_ylabel('Profit', color="#C0392B")
                ax1.set_xlabel('Date')
                plt.title("Revenue vs Profit Trend")
                plt.xticks(rotation=45)
                fig.tight_layout()
                visuals.append(("Revenue vs Profit Trend", self._buf_from_fig(fig)))
        except Exception:
            pass

        # 2) Profit Margin Trend
        try:
            if margin:
                date_info = self._first_date(df)
                if date_info:
                    dc, ser = date_info
                    tmp = df.copy()
                    tmp[dc] = pd.to_datetime(tmp[dc], errors='coerce')
                    tmp = tmp.dropna(subset=[dc])
                    fig, ax = plt.subplots(figsize=(10,5.2))
                    ax.plot(tmp[dc], tmp[margin]*100, color="#16A085", marker='o')
                    ax.axhline(y=(tmp[margin]*100).mean(), color="#8E44AD", linestyle='--', alpha=0.6, label='Average')
                    ax.set_title("Profit Margin (%) Over Time")
                    ax.set_xlabel("Date")
                    ax.set_ylabel("Profit Margin (%)")
                    ax.legend()
                    plt.xticks(rotation=45)
                    fig.tight_layout()
                    visuals.append(("Profit Margin Over Time", self._buf_from_fig(fig)))
        except Exception:
            pass

        # 3) Cost Breakdown stacked
        try:
            cost_cols = [c for c in df.columns if any(k in c.lower() for k in ['ingredients','packaging','labor','utilities','opex','operating_expenses','salaries','rent','marketing','maintenance','supplies','depreciation','cogs'])]
            if cost_cols:
                subset = df[cost_cols].select_dtypes(include=[np.number])
                if subset.shape[1]>=2:
                    totals = subset.sum().sort_values(ascending=False).head(8)
                    fig, ax = plt.subplots(figsize=(8.5,5.2))
                    bars = ax.bar(totals.index, totals.values, color=sns.color_palette("pastel", n_colors=len(totals)))
                    ax.set_title("Cost Breakdown (Total)")
                    ax.set_ylabel("Amount")
                    ax.set_xlabel("Cost Category")
                    plt.xticks(rotation=30, ha='right')
                    for b in bars:
                        h = b.get_height()
                        ax.text(b.get_x()+b.get_width()/2, h, f"{h:,.0f}", ha='center', va='bottom', fontsize=9)
                    fig.tight_layout()
                    visuals.append(("Cost Breakdown", self._buf_from_fig(fig)))
        except Exception:
            pass

        # 4) Correlation heatmap
        try:
            num = df.select_dtypes(include=[np.number])
            if num.shape[1]>=2:
                fig, ax = plt.subplots(figsize=(8.5,5.5))
                sns.heatmap(num.corr(), annot=True, fmt='.2f', cmap='coolwarm', center=0, ax=ax)
                ax.set_title("Correlation Matrix (Financial)")
                fig.tight_layout()
                visuals.append(("Correlation Matrix", self._buf_from_fig(fig)))
        except Exception:
            pass

        return visuals

    # Sales charts: top products, daily sales trend, payment method pie, store/category bars
    def _sales_charts(self, df):
        visuals = []
        cols = {c.lower(): c for c in df.columns}
        product = next((cols[c] for c in cols if 'product' in c), None)
        category = next((cols[c] for c in cols if 'category' in c), None)
        revenue = next((cols[c] for c in cols if c in ['total_sales','sales','revenue','amount','unit_price']), None)
        qty = next((cols[c] for c in cols if 'quantity' in c or 'qty' in c), None)
        payment = next((cols[c] for c in cols if 'payment' in c), None)

        # 1) Top Products by Revenue
        try:
            if product and revenue:
                top = df.groupby(product)[revenue].sum().sort_values(ascending=False).head(10)
                if top.shape[0]>0:
                    fig, ax = plt.subplots(figsize=(10,5.5))
                    bars = ax.bar(top.index, top.values, color=sns.color_palette("Blues_r", n_colors=len(top)))
                    ax.set_title("Top 10 Products by Revenue")
                    ax.set_xlabel("Product")
                    ax.set_ylabel("Revenue")
                    plt.xticks(rotation=40, ha='right')
                    for b in bars:
                        h = b.get_height()
                        ax.text(b.get_x()+b.get_width()/2, h, f"{h:,.0f}", ha='center', va='bottom', fontsize=9)
                    fig.tight_layout()
                    visuals.append(("Top Products by Revenue", self._buf_from_fig(fig)))
        except Exception:
            pass

        # 2) Daily Sales Trend
        try:
            date_info = self._first_date(df)
            if date_info and revenue:
                dc, ser = date_info
                tmp = df.copy()
                tmp[dc] = pd.to_datetime(tmp[dc], errors='coerce')
                tmp = tmp.dropna(subset=[dc])
                daily = tmp.groupby(tmp[dc].dt.to_period('D'))[revenue].sum()
                fig, ax = plt.subplots(figsize=(10,5.2))
                ax.plot(daily.index.astype(str), daily.values, marker='o', color="#2E86AB")
                z = np.polyfit(range(len(daily)), daily.values, 1)
                p = np.poly1d(z)
                ax.plot(daily.index.astype(str), p(range(len(daily))), "--", color='#C0392B', alpha=0.7)
                ax.set_title("Daily Sales Trend")
                ax.set_xlabel("Date")
                ax.set_ylabel("Sales")
                plt.xticks(rotation=45)
                fig.tight_layout()
                visuals.append(("Daily Sales Trend", self._buf_from_fig(fig)))
        except Exception:
            pass

        # 3) Payment Method Distribution
        try:
            if payment:
                dist = df[payment].value_counts().head(6)
                if dist.shape[0]>0:
                    fig, ax = plt.subplots(figsize=(7,7))
                    ax.pie(dist.values, labels=dist.index, autopct='%1.1f%%', startangle=90,
                           colors=sns.color_palette("pastel", n_colors=len(dist)))
                    ax.set_title("Payment Method Distribution")
                    visuals.append(("Payment Method Distribution", self._buf_from_fig(fig)))
        except Exception:
            pass

        # 4) Category performance
        try:
            if category and revenue:
                cat_perf = df.groupby(category)[revenue].sum().sort_values(ascending=False)
                if cat_perf.shape>0:
                    fig, ax = plt.subplots(figsize=(9.5,5.2))
                    bars = ax.bar(cat_perf.index, cat_perf.values, color=sns.color_palette("Greens", n_colors=len(cat_perf)))
                    ax.set_title("Revenue by Category")
                    ax.set_xlabel("Category")
                    ax.set_ylabel("Revenue")
                    plt.xticks(rotation=30, ha='right')
                    fig.tight_layout()
                    visuals.append(("Revenue by Category", self._buf_from_fig(fig)))
        except Exception:
            pass

        return visuals

    # Social media charts: platform average metric, content distribution, engagement correlation
    def _social_media_charts(self, df):
        visuals = []
        cols = {c.lower(): c for c in df.columns}
        platform = cols.get('platform')
        content = next((cols[c] for c in cols if 'content' in c), None)
        metrics = [c for c in df.columns if c.lower() in ['views','likes','shares','comments']]
        date_info = self._first_date(df)

        # 1) Average metric by platform
        try:
            if platform and metrics:
                metric = metrics[0]
                perf = df.groupby(platform)[metric].mean().sort_values(ascending=False)
                fig, ax = plt.subplots(figsize=(9.5,5.2))
                bars = ax.bar(perf.index, perf.values, color=sns.color_palette("coolwarm", n_colors=len(perf)))
                ax.set_title(f"Average {metric} by Platform")
                ax.set_xlabel("Platform")
                ax.set_ylabel(f"Avg {metric}")
                for b in bars:
                    h = b.get_height()
                    ax.text(b.get_x()+b.get_width()/2, h, f"{h:,.0f}", ha='center', va='bottom', fontsize=9)
                plt.xticks(rotation=20)
                fig.tight_layout()
                visuals.append((f"Average {metric} by Platform", self._buf_from_fig(fig)))
        except Exception:
            pass

        # 2) Content type distribution
        try:
            if content:
                cnt = df[content].value_counts().head(8)
                fig, ax = plt.subplots(figsize=(7.2,7.2))
                ax.pie(cnt.values, labels=cnt.index, autopct='%1.1f%%', startangle=90,
                       colors=sns.color_palette("pastel", n_colors=len(cnt)))
                ax.set_title("Content Type Distribution")
                visuals.append(("Content Type Distribution", self._buf_from_fig(fig)))
        except Exception:
            pass

        # 3) Engagement correlation
        try:
            eng = df[[c for c in metrics if df[c].dtype in [np.float64, np.int64, np.int32, np.float32]]]
            if eng.shape[1]>=2:
                fig, ax = plt.subplots(figsize=(8.5,5.5))
                sns.heatmap(eng.corr(), annot=True, fmt='.2f', cmap='coolwarm', center=0, ax=ax)
                ax.set_title("Engagement Metrics Correlation")
                fig.tight_layout()
                visuals.append(("Engagement Metrics Correlation", self._buf_from_fig(fig)))
        except Exception:
            pass

        # 4) Time trend on primary metric
        try:
            if date_info and metrics:
                dc, ser = date_info
                tmp = df.copy()
                tmp[dc] = pd.to_datetime(tmp[dc], errors='coerce')
                tmp = tmp.dropna(subset=[dc])
                daily = tmp.groupby(tmp[dc].dt.to_period('D'))[metrics].sum()
                fig, ax = plt.subplots(figsize=(10,5.2))
                ax.plot(daily.index.astype(str), daily.values, color="#2E86AB", marker='o')
                ax.set_title(f"Daily {metrics} Trend")
                ax.set_xlabel("Date")
                ax.set_ylabel(metrics)
                plt.xticks(rotation=45)
                fig.tight_layout()
                visuals.append((f"Daily {metrics} Trend", self._buf_from_fig(fig)))
        except Exception:
            pass

        return visuals

    # General charts: numeric trend, top category by numeric, distribution pie, correlation
    def _general_charts(self, df):
        visuals = []
        num_cols = df.select_dtypes(include=[np.number]).columns.tolist()
        txt_cols = df.select_dtypes(include=['object']).columns.tolist()
        date_info = self._first_date(df)

        # 1) Numeric trend over time
        try:
            if date_info and num_cols:
                dc, ser = date_info
                tmp = df.copy()
                tmp[dc] = pd.to_datetime(tmp[dc], errors='coerce')
                tmp = tmp.dropna(subset=[dc])
                daily = tmp.groupby(tmp[dc].dt.to_period('D'))[num_cols].sum()
                fig, ax = plt.subplots(figsize=(10,5.2))
                ax.plot(daily.index.astype(str), daily.values, color="#2E86AB", marker='o')
                ax.set_title(f"{num_cols.title()} Trend Over Time")
                ax.set_xlabel("Date")
                ax.set_ylabel(num_cols.title())
                plt.xticks(rotation=45)
                fig.tight_layout()
                visuals.append((f"{num_cols.title()} Trend", self._buf_from_fig(fig)))
        except Exception:
            pass

        # 2) Top categories by first numeric
        try:
            if txt_cols and num_cols:
                top = df.groupby(txt_cols)[num_cols].sum().nlargest(10)
                fig, ax = plt.subplots(figsize=(10,5.2))
                bars = ax.bar(range(len(top)), top.values, color=sns.color_palette("muted", n_colors=len(top)))
                ax.set_title(f"Top {txt_cols.title()} by {num_cols.title()}")
                ax.set_xlabel(txt_cols.title())
                ax.set_ylabel(num_cols.title())
                ax.set_xticks(range(len(top)))
                ax.set_xticklabels(top.index, rotation=35, ha='right')
                for b in bars:
                    h = b.get_height()
                    ax.text(b.get_x()+b.get_width()/2, h, f"{h:,.0f}", ha='center', va='bottom', fontsize=9)
                fig.tight_layout()
                visuals.append((f"Top {txt_cols.title()} by {num_cols.title()}", self._buf_from_fig(fig)))
        except Exception:
            pass

        # 3) Pie distribution
        try:
            if txt_cols:
                dist = df[txt_cols].value_counts().head(8)
                fig, ax = plt.subplots(figsize=(7.2,7.2))
                ax.pie(dist.values, labels=dist.index, autopct='%1.1f%%', startangle=90,
                       colors=sns.color_palette("pastel", n_colors=len(dist)))
                ax.set_title(f"{txt_cols.title()} Distribution")
                visuals.append((f"{txt_cols.title()} Distribution", self._buf_from_fig(fig)))
        except Exception:
            pass

        # 4) Correlation matrix
        try:
            if len(num_cols) >= 2:
                fig, ax = plt.subplots(figsize=(8.5,5.5))
                sns.heatmap(df[num_cols].corr(), annot=True, fmt='.2f', cmap='coolwarm', center=0, ax=ax)
                ax.set_title("Correlation Matrix")
                fig.tight_layout()
                visuals.append(("Correlation Matrix", self._buf_from_fig(fig)))
        except Exception:
            pass

        return visuals

    def create_feedback_report(self, df, column_analysis, feedback_analysis, visualizations, executive_summary, filename):
        """Generate comprehensive PDF report with AI-enhanced insights in markdown style"""
        try:
            from reportlab.lib.pagesizes import A4, letter
            from reportlab.lib import colors
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image
            from reportlab.lib.units import inch
            from io import BytesIO
            from datetime import datetime
            import requests
            
            # Create PDF buffer
            buffer = BytesIO()
            doc = SimpleDocTemplate(buffer, pagesize=A4, 
                                topMargin=0.5*inch, 
                                bottomMargin=0.5*inch,
                                leftMargin=0.5*inch,
                                rightMargin=0.5*inch)
            
            styles = getSampleStyleSheet()
            story = []

            # Custom styles for AI-enhanced report
            title_style = ParagraphStyle(
                'AITitle',
                parent=styles['Heading1'],
                fontSize=22,
                spaceAfter=20,
                textColor=colors.HexColor('#2C3E50'),
                alignment=1,
                fontName='Helvetica-Bold'
            )

            section_style = ParagraphStyle(
                'AISection',
                parent=styles['Heading2'],
                fontSize=16,
                spaceAfter=12,
                spaceBefore=15,
                textColor=colors.HexColor("#090909"),
                fontName='Helvetica-Bold'
            )

            subsection_style = ParagraphStyle(
                'AISubsection',
                parent=styles['Heading3'],
                fontSize=13,
                spaceAfter=8,
                spaceBefore=10,
                textColor=colors.HexColor('#4c5c68'),
                fontName='Helvetica-Bold'
            )

            bullet_style = ParagraphStyle(
                'AIBullet',
                parent=styles['Normal'],
                fontSize=11,
                spaceAfter=4,
                leftIndent=20,
                bulletIndent=10,
                textColor=colors.HexColor("#000000"),
                fontName='Helvetica'
            )

            insight_style = ParagraphStyle(
                'AIInsight',
                parent=styles['Normal'],
                fontSize=11,
                spaceAfter=4,
                leftIndent=30,
                bulletIndent=15,
                textColor=colors.HexColor("#000000"),
                fontName='Helvetica'
            )

            recommendation_style = ParagraphStyle(
                'AIRecommendation',
                parent=styles['Normal'],
                fontSize=11,
                spaceAfter=4,
                leftIndent=30,
                bulletIndent=15,
                textColor=colors.HexColor("#000000"),
                fontName='Helvetica'
            )

            quote_style = ParagraphStyle(
                'AIQuote',
                parent=styles['Italic'],
                fontSize=10,
                spaceAfter=4,
                leftIndent=40,
                textColor=colors.HexColor("#000000"),
                fontName='Helvetica-Oblique'
            )

            highlight_style = ParagraphStyle(
                'AIHighlight',
                parent=styles['Normal'],
                fontSize=11,
                spaceAfter=4,
                leftIndent=20,
                textColor=colors.HexColor('#E67E22'),
                fontName='Helvetica-Bold'
            )

            # Title and Executive Summary
            story.append(Paragraph("AI-Powered Customer Feedback Analysis Report", title_style))
            story.append(Spacer(1, 15))
            
            # Executive Summary from AI analysis - Parse markdown format
            story.append(Paragraph("Executive Summary", section_style))
            
            if executive_summary:
                self.parse_markdown_executive_summary(executive_summary, story, subsection_style, bullet_style, insight_style, recommendation_style, highlight_style)
            else:
                story.append(Paragraph("• Comprehensive analysis of customer feedback revealing key insights", bullet_style))
                story.append(Paragraph("• Identification of improvement opportunities", bullet_style))
            
            story.append(Spacer(1, 20))

            # Metadata in bullet format
            story.append(Paragraph("Report Overview", section_style))
            story.append(Paragraph(f"• File Name: {filename}", bullet_style))
            story.append(Paragraph(f"• Total Records: {len(df):,}", bullet_style))
            story.append(Paragraph(f"• Analysis Date: {datetime.now().strftime('%Y-%m-%d %H:%M')}", bullet_style))
            story.append(Paragraph(f"• Data Quality Score: {(df.notna().sum().sum() / (df.shape[0] * df.shape[1]) * 100):.1f}%", bullet_style))
            story.append(Paragraph(f"• Columns Analyzed: {len([col for col in column_analysis.get('columns_analysis', []) if col.get('include_in_analysis', True)])}", bullet_style))
            story.append(Spacer(1, 25))

            # Sentiment Analysis Section - Bullet Format
            sentiment_data = feedback_analysis.get('sentiment_summary', {})
            story.append(Paragraph("Comprehensive Sentiment Analysis", section_style))
            
            if sentiment_data:
                story.append(Paragraph("Sentiment Distribution:", subsection_style))
                story.append(Paragraph(f"• Positive: {sentiment_data.get('positive_percentage', 0):.1f}% - Satisfaction, praise, positive experiences", insight_style))
                story.append(Paragraph(f"• Negative: {sentiment_data.get('negative_percentage', 0):.1f}% - Issues, complaints, dissatisfaction", recommendation_style))
                story.append(Paragraph(f"• Neutral: {sentiment_data.get('neutral_percentage', 0):.1f}% - Factual information, suggestions, mixed feedback", bullet_style))
                story.append(Paragraph(f"• Overall Sentiment: {sentiment_data.get('overall_sentiment', 'N/A').title()}", highlight_style))
                story.append(Spacer(1, 10))

                # Example quotes from each sentiment
                sentiment_examples = sentiment_data.get('sentiment_distribution', {})
                if sentiment_examples:
                    story.append(Paragraph("Representative Feedback Examples:", subsection_style))
                    
                    for sentiment, examples in sentiment_examples.items():
                        if examples and len(examples) > 0:
                            story.append(Paragraph(f"{sentiment.title()} Examples:", bullet_style))
                            for example in examples[:2]:  # Show 2 examples per sentiment
                                story.append(Paragraph(f'"{example}"', quote_style))
                            story.append(Spacer(1, 5))
            else:
                story.append(Paragraph("• No sentiment data available", bullet_style))
            
            story.append(Spacer(1, 20))

            # Positive Feedback Analysis - Bullet Format
            positive_analysis = feedback_analysis.get('positive_feedback_analysis', {})
            story.append(Paragraph("Positive Feedback Analysis", section_style))
            
            if positive_analysis:
                positive_cats = positive_analysis.get('categories', [])
                
                if positive_cats:
                    story.append(Paragraph("Top Positive Categories:", subsection_style))
                    for cat in positive_cats[:5]:  # Show top 5
                        category_name = cat.get('category', 'N/A')
                        percentage = cat.get('percentage', 0)
                        themes = ', '.join(cat.get('key_themes', [])[:3])  # Show top 3 themes
                        story.append(Paragraph(f"• {category_name}: {percentage:.1f}%", insight_style))
                        if themes:
                            story.append(Paragraph(f"  - Key themes: {themes}", bullet_style))
                    story.append(Spacer(1, 10))

                # Top positive aspects
                top_positives = positive_analysis.get('top_positive_aspects', [])
                if top_positives:
                    story.append(Paragraph("Top Positive Aspects:", subsection_style))
                    for aspect in top_positives[:5]:
                        story.append(Paragraph(f"• {aspect}", insight_style))
            else:
                story.append(Paragraph("• No positive feedback analysis available", bullet_style))
            
            story.append(Spacer(1, 20))

            # Negative Feedback Analysis - Bullet Format
            negative_analysis = feedback_analysis.get('negative_feedback_analysis', {})
            story.append(Paragraph("Negative Feedback Analysis", section_style))
            
            if negative_analysis:
                negative_cats = negative_analysis.get('categories', [])
                
                if negative_cats:
                    story.append(Paragraph("Top Negative Categories:", subsection_style))
                    for cat in negative_cats[:5]:  # Show top 5
                        category_name = cat.get('category', 'N/A')
                        percentage = cat.get('percentage', 0)
                        issues = ', '.join(cat.get('key_issues', [])[:3])  # Show top 3 issues
                        story.append(Paragraph(f"• {category_name}: {percentage:.1f}%", recommendation_style))
                        if issues:
                            story.append(Paragraph(f"  - Main issues: {issues}", bullet_style))
                    story.append(Spacer(1, 10))

                # Top improvement areas
                top_improvements = negative_analysis.get('top_improvement_areas', [])
                if top_improvements:
                    story.append(Paragraph("Top Improvement Areas:", subsection_style))
                    for area in top_improvements[:5]:
                        story.append(Paragraph(f"• {area}", recommendation_style))
            else:
                story.append(Paragraph("• No negative feedback analysis available", bullet_style))
            
            story.append(Spacer(1, 20))

            # Actionable Recommendations - Bullet Format
            recommendations = feedback_analysis.get('recommendations', [])
            story.append(Paragraph("AI-Generated Recommendations", section_style))
            
            if recommendations:
                # Group recommendations by priority
                high_priority = [r for r in recommendations if r.get('priority', '').lower() == 'high']
                medium_priority = [r for r in recommendations if r.get('priority', '').lower() == 'medium']
                low_priority = [r for r in recommendations if r.get('priority', '').lower() == 'low']
                
                if high_priority:
                    story.append(Paragraph("High Priority Actions:", subsection_style))
                    for rec in high_priority[:5]:
                        area = rec.get('area', 'General')
                        action = rec.get('action', 'No action specified')
                        timeline = rec.get('timeline', 'Short-term')
                        story.append(Paragraph(f"• {area}: {action}", recommendation_style))
                        story.append(Paragraph(f"  - Timeline: {timeline}", bullet_style))
                    story.append(Spacer(1, 10))
                
                if medium_priority:
                    story.append(Paragraph("Medium Priority Actions:", subsection_style))
                    for rec in medium_priority[:3]:
                        area = rec.get('area', 'General')
                        action = rec.get('action', 'No action specified')
                        timeline = rec.get('timeline', 'Medium-term')
                        story.append(Paragraph(f"• {area}: {action}", bullet_style))
                        story.append(Paragraph(f"  - Timeline: {timeline}", bullet_style))
                    story.append(Spacer(1, 10))
                
                if low_priority:
                    story.append(Paragraph("Low Priority Actions:", subsection_style))
                    for rec in low_priority[:3]:
                        area = rec.get('area', 'General')
                        action = rec.get('action', 'No action specified')
                        story.append(Paragraph(f"• {area}: {action}", bullet_style))
            else:
                story.append(Paragraph("• No specific recommendations available", bullet_style))
            
            story.append(Spacer(1, 25))

            # Data Visualizations Section
            story.append(Paragraph("AI-Enhanced Visualizations", section_style))
            story.append(Spacer(1, 10))

            if visualizations:
                # Display each visualization with AI-generated description
                for i, viz in enumerate(visualizations, 1):
                    try:
                        story.append(Paragraph(f"Chart {i}: {viz.get('title', 'Visualization')}", subsection_style))
                        
                        # Download the chart image from Supabase URL
                        chart_url = viz.get('url')
                        if chart_url:
                            response = requests.get(chart_url, timeout=10)
                            if response.status_code == 200:
                                chart_buffer = BytesIO(response.content)
                                
                                # Convert to ReportLab Image with appropriate sizing
                                chart_image = Image(chart_buffer, width=5.5*inch, height=3.5*inch)
                                chart_image.hAlign = 'CENTER'
                                story.append(chart_image)
                                
                                # Add AI-generated chart description as bullet points
                                description = viz.get('description', 'No description available')
                                # Split description into key points
                                desc_points = description.split('. ')
                                story.append(Paragraph("📝 Key Insights:", bullet_style))
                                for point in desc_points[:3]:  # Show top 3 insights
                                    if point.strip():
                                        clean_point = point.strip().rstrip('.')
                                        story.append(Paragraph(f"• {clean_point}", insight_style))
                                story.append(Spacer(1, 15))
                    except Exception as e:
                        print(f"Error loading chart {i}: {e}")
                        continue
            else:
                story.append(Paragraph("• No visualizations available for this analysis", bullet_style))
            
            story.append(Spacer(1, 25))

            # Column Analysis Summary - Bullet Format
            # story.append(Paragraph("Data Structure Insights", section_style))
            
            # columns_analysis = column_analysis.get('columns_analysis', [])
            # relevant_columns = [col for col in columns_analysis if col.get('include_in_analysis', True)]
            
            # if relevant_columns:
            #     story.append(Paragraph("Key Analyzed Columns:", subsection_style))
                
            #     for col_analysis in relevant_columns[:8]:  # Show top 8 relevant columns
            #         col_name = col_analysis.get('column_name', 'N/A')
            #         col_type = col_analysis.get('detected_type', 'unknown')
            #         reasoning = col_analysis.get('reasoning', 'No purpose specified')[:50] + '...'
                    
            #         story.append(Paragraph(f"• {col_name} ({col_type})", bullet_style))
            #         story.append(Paragraph(f"  - Purpose: {reasoning}", bullet_style))
            #         story.append(Spacer(1, 2))
            # else:
            #     story.append(Paragraph("• No column analysis available", bullet_style))
            
            # story.append(Spacer(1, 25))

            # Methodology and Limitations - Bullet Format
            story.append(Paragraph("Analysis Methodology", section_style))
            story.append(Paragraph("Advanced AI Techniques Used:", subsection_style))
            story.append(Paragraph("• Natural Language Processing for sentiment analysis and theme extraction", bullet_style))
            story.append(Paragraph("• Machine Learning for pattern recognition and categorization", bullet_style))
            story.append(Paragraph("• Statistical analysis for quantitative insights", bullet_style))
            story.append(Paragraph("• AI-powered visualization generation and interpretation", bullet_style))
            story.append(Spacer(1, 10))
            
            story.append(Paragraph("Important Limitations:", subsection_style))
            story.append(Paragraph("• Analysis is based on available data and may not capture all nuances", bullet_style))
            story.append(Paragraph("• Recommendations should be validated with additional business context", bullet_style))
            story.append(Paragraph("• AI confidence varies based on data quality and volume", bullet_style))
            story.append(Spacer(1, 20))

            # Footer
            footer_style = ParagraphStyle(
                'AIFooter',
                parent=styles['Normal'],
                fontSize=9,
                textColor=colors.HexColor('#7F8C8D'),
                alignment=1,
                fontName='Helvetica-Oblique'
            )
            
            story.append(Paragraph("Generated by Advanced AI Feedback Analysis System", footer_style))
            story.append(Paragraph(f"Report generated on: {datetime.now().strftime('%B %d, %Y at %I:%M %p')}", footer_style))
            story.append(Paragraph("Confidential - For strategic decision making only", footer_style))
            story.append(Paragraph("AI Confidence Level: High (Based on comprehensive pattern recognition)", footer_style))

            # Build the PDF
            doc.build(story)
            pdf_content = buffer.getvalue()
            buffer.close()
            
            return pdf_content
        
        except Exception as e:
            print(f"AI Enhanced PDF Generation Error: {str(e)}")
            import traceback
            traceback.print_exc()
            raise Exception(f"Failed to generate AI-enhanced PDF: {str(e)}")
    
    def parse_markdown_executive_summary(self, executive_summary, story, subsection_style, bullet_style, insight_style, recommendation_style, highlight_style):
        """Parse markdown formatted executive summary and convert to ReportLab format"""
        from reportlab.lib.styles import ParagraphStyle
        from reportlab.lib import colors
        from reportlab.platypus import Paragraph, Spacer
        
        # Create bold style for **text**
        bold_style = ParagraphStyle(
            'Bold',
            parent=bullet_style,
            fontName='Helvetica-Bold'
        )
        title, h2, h3, normal, bullet, note = _styles()
        lines = executive_summary.split('\n')
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
                
            # Handle headers (##, ###, etc.)
            if line.startswith('## '):
                header_text = line.replace('## ', '').replace('■ ', '').strip()
                story.append(Spacer(1, 6))
                story.append(Paragraph(f"{header_text}", subsection_style))
            elif line.startswith('###'):
                sub_style = ParagraphStyle(
                    'SubH3', parent=normal, fontName='Helvetica-Bold',
                    fontSize=12, textColor=colors.HexColor("#9B2C2C"), spaceBefore=8, spaceAfter=6
                )
                story.append(Paragraph(line[4:].strip(), sub_style)) 
            # Handle main bullets (*)
            elif line.startswith('* ') and not line.startswith('** '):
                bullet_text = line.replace('* ', '').strip()
                # Handle bold text within bullets
                bullet_text = self.format_bold_text(bullet_text)
                
                # Use different styles based on content
                if 'High Priority' in bullet_text or 'Medium Priority' in bullet_text:
                    story.append(Paragraph(f"• {bullet_text}", recommendation_style))
                elif any(word in bullet_text.lower() for word in ['positive', 'excellent', 'high-quality', 'praised']):
                    story.append(Paragraph(f"• {bullet_text}", insight_style))
                else:
                    story.append(Paragraph(f"• {bullet_text}", bullet_style))
                    
            # Handle numbered lists
            elif line[0].isdigit() and line[1:3] == '. ':
                numbered_text = line[3:].strip()
                numbered_text = self.format_bold_text(numbered_text)
                story.append(Paragraph(f"{line[:2]} {numbered_text}", bullet_style))
                
            # Handle regular paragraphs (not starting with * or #)
            elif not line.startswith(('# ', '* ')):
                formatted_text = self.format_bold_text(line)
                story.append(Paragraph(formatted_text, bullet_style))
                story.append(Spacer(1, 3))
    
    def format_bold_text(self, text):
        """Convert **text** to <b>text</b> for ReportLab"""
        import re
        # Replace **text** with <b>text</b>
        text = re.sub(r'\*\*(.*?)\*\*', r'<b>\1</b>', text)
        return text
              


# ---------------------------
# PPT Generator
# ---------------------------


class PPTGenerator:
    def create_analysis_presentation(self, gemini_output, filename, data_df=None, data_type=None):
        if data_df is not None and isinstance(data_df, pd.DataFrame) and not data_df.empty:
            return self.create_analysis_presentation_with_charts(gemini_output, data_df, filename, data_type)

        prs = self._new_presentation_theme()
        if isinstance(gemini_output, dict) and 'full_report' in gemini_output:
            report_content = gemini_output['full_report']
            chart_recommendations = gemini_output.get('chart_recommendations', [])
        else:
            report_content = str(gemini_output)
            chart_recommendations = []

        sections = self.parse_report_sections(report_content)

        self.create_title_slide(prs, sections.get('title', self._title_for_type(filename, data_type)))
        if 'executive_summary' in sections:
            self.create_content_slide(prs, "Executive Summary", sections['executive_summary'])
        if 'key_performance_indicators' in sections:
            self.create_content_slide(prs, "Key Performance Indicators", sections['key_performance_indicators'])
        if 'trend_analysis' in sections:
            self.create_content_slide(prs, "Trend Analysis", sections['trend_analysis'])
        if 'product_analysis' in sections:
            self.create_content_slide(prs, "Product Mix Analysis", sections['product_analysis'])
        if 'customer_insights' in sections:
            self.create_content_slide(prs, "Customer Behavior Insights", sections['customer_insights'])
        if 'recommendations' in sections:
            self.create_recommendations_slide(prs, sections['recommendations'])
        if chart_recommendations:
            self.create_visualization_recommendations_slide(prs, chart_recommendations)
        if 'conclusion' in sections:
            self.create_content_slide(prs, "Conclusion & Strategic Outlook", sections['conclusion'])

        buf = BytesIO()
        prs.save(buf)
        return buf.getvalue()

    def create_analysis_presentation_with_charts(self, gemini_output, data_df, filename, data_type=None):
        prs = self._new_presentation_theme()

        if isinstance(gemini_output, dict) and 'full_report' in gemini_output:
            report_content = gemini_output['full_report']
        else:
            report_content = str(gemini_output)

        sections = self.parse_report_sections(report_content)

        self.create_title_slide(prs, sections.get('title', self._title_for_type(filename, data_type)))
        if 'executive_summary' in sections:
            self.create_content_slide(prs, "Executive Summary", sections['executive_summary'])

        # Chart slides: data-type aware
        self.create_chart_slides(prs, data_df, data_type)

        if 'recommendations' in sections:
            self.create_recommendations_slide(prs, sections['recommendations'])
        if 'conclusion' in sections:
            self.create_content_slide(prs, "Conclusion & Strategic Outlook", sections['conclusion'])

        buf = BytesIO()
        prs.save(buf)
        return buf.getvalue()

    def create_specialized_analysis_presentation(self, analysis_results, filename, df_cleaned, data_type):
        """Create a specialized analysis presentation based on data type"""
        return self.create_analysis_presentation(analysis_results, filename, df_cleaned, data_type)

    def _new_presentation_theme(self):
        prs = Presentation()
        # Basic theme styling can be adjusted per brand
        return prs

    def _title_for_type(self, filename, data_type):
        base = "Business Analysis Report"
        if data_type == 'sales':
            base = "Sales Performance Report"
        elif data_type == 'financial':
            base = "Financial Performance Report"
        elif data_type == 'social_media':
            base = "Social Media Performance Report"
        return f"{base} – {filename}"

    def parse_report_sections(self, report_content):
        sections = {}
        lines = report_content.split('\n')
        current = ""
        buf = []
        for line in lines:
            line = line.strip()
            if line.startswith('# '):
                sections['title'] = line[2:].strip()
            elif line.startswith('## '):
                if current and buf:
                    sections[current] = '\n'.join(buf)
                current = line[3:].strip().lower().replace(' ', '_').replace('&amp;', 'and').replace('&', 'and')
                buf = []
            elif line:
                buf.append(line)
        if current and buf:
            sections[current] = '\n'.join(buf)
        return sections

    def create_title_slide(self, prs, title):
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
        subtitle = slide.shapes.placeholders[1]
        subtitle.text = f"Executive Business Intelligence Report\n{datetime.datetime.now().strftime('%B %Y')}\nGenerated by AI Analytics"

        # Style
        title_tf = slide.shapes.title.text_frame
        title_tf.paragraphs[0].font.size = Pt(36)
        title_tf.paragraphs[0].font.bold = True
        title_tf.paragraphs[0].font.color.theme_color = MSO_THEME_COLOR.ACCENT_1

        stf = subtitle.text_frame
        for p in stf.paragraphs:
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(80, 80, 80)

    def create_content_slide(self, prs, title, content):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title

        # Title style
        title_tf = slide.shapes.title.text_frame
        title_tf.paragraphs.font.size = Pt(28)
        title_tf.paragraphs.font.bold = True
        title_tf.paragraphs.font.color.rgb = RGBColor(20, 61, 107)

        tf = slide.shapes.placeholders[1].text_frame
        tf.clear()
        lines = content.split('\n')
        first = True
        for line in lines:
            line = line.strip()
            if not line:
                continue
            para = tf.paragraphs if first else tf.add_paragraph()
            first = False
            if line.startswith('###'):
                para.text = line[3:].strip()
                para.level = 0
                para.font.bold = True
                para.font.size = Pt(18)
                para.font.color.rgb = RGBColor(155, 44, 44)
            elif line.startswith(('- ', '* ')) or line[:3] in [f"{i}. " for i in range(1,10)]:
                para.text = line[2:] if line.startswith(('- ', '* ')) else line
                para.level = 1
                para.font.size = Pt(16)
            elif any(emoji in line for emoji in ['📊','📈','🎯']):
                para.text = f"• {line.replace('📊','').replace('📈','').replace('🎯','').strip()}"
                para.level = 1
                para.font.size = Pt(14)
                para.font.color.rgb = RGBColor(30, 158, 106)
            else:
                para.text = line
                para.level = 1
                para.font.size = Pt(16)

    def create_recommendations_slide(self, prs, recommendations_content):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Strategic Recommendations & Action Plan"

        title_tf = slide.shapes.title.text_frame
        title_tf.paragraphs.font.size = Pt(28)
        title_tf.paragraphs.font.bold = True
        title_tf.paragraphs.font.color.rgb = RGBColor(20, 61, 107)

        tf = slide.shapes.placeholders[1].text_frame
        tf.clear()

        priority_colors = {
            'immediate': RGBColor(197, 48, 48),
            'strategic': RGBColor(221, 107, 32),
            'long': RGBColor(0, 100, 0),
        }

        current_priority = None
        first = True
        for line in recommendations_content.split('\n'):
            line = line.strip()
            if not line:
                continue
            para = tf.paragraphs if first else tf.add_paragraph()
            first = False
            lower = line.lower()
            if 'immediate' in lower or 'next 30' in lower:
                current_priority = 'immediate'
                para.text = "IMMEDIATE ACTIONS (Next 30 Days)"
                para.level = 0
                para.font.size = Pt(18)
                para.font.bold = True
                para.font.color.rgb = priority_colors['immediate']
            elif 'strategic' in lower or 'quarter' in lower:
                current_priority = 'strategic'
                para.text = "STRATEGIC INITIATIVES (Next Quarter)"
                para.level = 0
                para.font.size = Pt(18)
                para.font.bold = True
                para.font.color.rgb = priority_colors['strategic']
            elif 'long' in lower or 'year' in lower:
                current_priority = 'long'
                para.text = "LONG-TERM GROWTH (Next Year)"
                para.level = 0
                para.font.size = Pt(18)
                para.font.bold = True
                para.font.color.rgb = priority_colors['long']
            elif line.startswith(('- ', '* ')) or line[:3] in [f"{i}. " for i in range(1,10)]:
                para.text = line[2:] if line.startswith(('- ', '* ')) else line
                para.level = 1
                para.font.size = Pt(16)
                if current_priority:
                    para.font.color.rgb = priority_colors[current_priority]
            else:
                para.text = line
                para.level = 1
                para.font.size = Pt(16)

    def create_visualization_recommendations_slide(self, prs, chart_recommendations):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Recommended Data Visualizations"
        title_tf = slide.shapes.title.text_frame
        title_tf.paragraphs.font.size = Pt(28)
        title_tf.paragraphs.font.bold = True
        title_tf.paragraphs.font.color.rgb = RGBColor(20, 61, 107)

        tf = slide.shapes.placeholders[1].text_frame
        tf.text = "Key Charts & Graphs"
        tf.paragraphs.font.size = Pt(16)
        for chart_rec in chart_recommendations:
            clean = chart_rec.replace('📊','').replace('📈','').replace('🎯','').strip()
            p = tf.add_paragraph()
            p.text = f"• {clean}"
            p.level = 1
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(30, 158, 106)

    def create_chart_slides(self, prs, df, data_type=None):
        # Use native pptx charts where feasible; otherwise embed rendered images
        # We’ll create 3-4 slides tailored by data type
        cols = {c.lower(): c for c in df.columns}

        def add_line_chart(title, categories, series_label, series_values):
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = title
            chart_data = CategoryChartData()
            chart_data.categories = categories
            chart_data.add_series(series_label, series_values)
            x, y, cx, cy = Inches(1), Inches(1.8), Inches(8), Inches(4.5)
            chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, x, y, cx, cy, chart_data).chart
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM

        def add_bar_chart(title, cats, label, vals):
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = title
            chart_data = CategoryChartData()
            chart_data.categories = cats
            chart_data.add_series(label, vals)
            x, y, cx, cy = Inches(1), Inches(1.8), Inches(8), Inches(4.5)
            chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart
            chart.has_legend = False

        def add_pie_chart(title, cats, vals):
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = title
            chart_data = CategoryChartData()
            chart_data.categories = cats
            chart_data.add_series('Share', vals)
            x, y, cx, cy = Inches(2), Inches(1.8), Inches(6), Inches(4.5)
            chart = slide.shapes.add_chart(XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data).chart
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.RIGHT

        # Date handling
        date_col = next((cols[c] for c in cols if 'date' in c or 'time' in c), None)
        date_series = None
        if date_col is not None:
            try:
                date_series = pd.to_datetime(df[date_col], errors='coerce')
            except Exception:
                date_series = None

        # Sales
        if data_type == 'sales':
            revenue = next((x for x in [cols.get('total_sales'), cols.get('revenue'), cols.get('sales'), cols.get('amount')] if x is not None), None)
            product = next((x for x in [cols.get('product'), cols.get('category'), cols.get('item')] if x is not None), None)
            payment = next((x for x in [cols.get('payment_method'), cols.get('payment_type'), cols.get('method')] if x is not None), None)

            # Trend
            if date_series is not None and revenue is not None:
                tmp = df.copy()
                tmp[date_col] = date_series
                tmp = tmp.dropna(subset=[date_col])
                daily = tmp.groupby(tmp[date_col].dt.to_period('D'))[revenue].sum()
                if len(daily)>0:
                    add_line_chart("Revenue Trend", [str(p) for p in daily.index], "Revenue", daily.values)

            # Top products
            if product and revenue:
                top = df.groupby(product)[revenue].sum().nlargest(5)
                if len(top)>0:
                    add_bar_chart("Top Performing Products", list(top.index), "Revenue", list(top.values))

            # Payment pie
            if payment:
                dist = df[payment].value_counts().head(6)
                if len(dist)>0:
                    add_pie_chart("Payment Method Distribution", list(dist.index), list(dist.values))
            return

        # Financial
        if data_type == 'financial':
            revenue = next((x for x in [cols.get('revenue'), cols.get('total_revenue')] if x is not None), None)
            profit = next((x for x in [cols.get('net_profit'), cols.get('profit'), cols.get('gross_profit')] if x is not None), None)
            margin = next((x for x in [cols.get('profit_margin'), cols.get('margin')] if x is not None), None)

            if date_series is not None and revenue and profit:
                tmp = df.copy()
                tmp[date_col] = date_series
                tmp = tmp.dropna(subset=[date_col])
                daily = tmp.groupby(tmp[date_col].dt.to_period('D'))[[revenue, profit]].sum()
                if len(daily)>0:
                    # Revenue line
                    add_line_chart("Revenue Trend", [str(p) for p in daily.index], "Revenue", daily[revenue].values)
                    # Profit line
                    add_line_chart("Profit Trend", [str(p) for p in daily.index], "Profit", daily[profit].values)

            if date_series is not None and margin:
                tmp = df.copy()
                tmp[date_col] = date_series
                tmp = tmp.dropna(subset=[date_col])
                series = tmp.groupby(tmp[date_col].dt.to_period('D'))[margin].mean()*100
                if len(series)>0:
                    add_line_chart("Profit Margin Trend (%)", [str(p) for p in series.index], "Margin %", series.values)
            return

        # Social Media
        if data_type == 'social_media':
            platform = cols.get('platform')
            metric = next((m for m in ['views','likes','shares','comments'] if m in cols), None)
            
            # Fix the content column lookup
            content = None
            for col in df.columns:
                if 'content' in col.lower():
                    content = col  
                    break
            
            if platform and metric:
                perf = df.groupby(platform)[cols[metric]].mean().sort_values(ascending=False).head(6)
                if len(perf) > 0:
                    add_bar_chart(f"Average {metric.title()} by Platform", list(perf.index), metric.title(), list(perf.values))
            
            if content:
                cnt = df[content].value_counts().head(6)
                if len(cnt) > 0:
                    add_pie_chart("Content Type Distribution", list(cnt.index), list(cnt.values))

            if date_series is not None and metric:
                tmp = df.copy()
                tmp[date_col] = date_series
                tmp = tmp.dropna(subset=[date_col])
                daily = tmp.groupby(tmp[date_col].dt.to_period('D'))[cols[metric]].sum()
                if len(daily)>0:
                    add_line_chart(f"Daily {metric.title()} Trend", [str(p) for p in daily.index], metric.title(), daily.values)
            return

        # General fallback
        num_cols = df.select_dtypes(include=[np.number]).columns.tolist()
        txt_cols = df.select_dtypes(include=['object']).columns.tolist()

        if date_series is not None and num_cols:
            tmp = df.copy()
            tmp[date_col] = date_series
            tmp = tmp.dropna(subset=[date_col])
            daily = tmp.groupby(tmp[date_col].dt.to_period('D'))[num_cols[0]].sum()
            if len(daily)>0:
                add_line_chart(f"{num_cols.title()} Trend", [str(p) for p in daily.index], num_cols.title(), daily.values)

        if txt_cols and num_cols:
            top = df.groupby(txt_cols)[num_cols].sum().nlargest(5)
            if len(top)>0:
                add_bar_chart(f"Top {txt_cols.title()}", list(top.index), num_cols.title(), list(top.values))


# End of file
